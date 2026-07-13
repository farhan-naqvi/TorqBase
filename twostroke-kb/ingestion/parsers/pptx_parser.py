"""PowerPoint (.pptx) -> slide text + speaker notes + embedded tables."""
from __future__ import annotations

import logging
import re
from pathlib import Path

from ..types import ParsedDoc, Table
from .math_normalizer import normalize_math_unicode

log = logging.getLogger(__name__)
_MIN_IMAGE_BYTES = 2048
_MAX_IMAGE_BYTES = 10 * 1024 * 1024
_CAPTION_RE = re.compile(r"\b(Figure|Fig\.|Abb\.|Bild|Diagram)\b", re.IGNORECASE)


def _doc_id_slug(path: Path) -> str:
    return re.sub(r"[^a-z0-9_.]", "_", path.name.lower())


def _image_filename_slug(path: Path) -> str:
    return re.sub(r"[^a-z0-9_.-]", "_", path.stem.lower())


def _already_extracted(doc_id: str) -> bool:
    try:
        from config import get_connection

        conn = get_connection()
        try:
            cur = conn.cursor()
            cur.execute("SELECT 1 FROM document_images WHERE doc_id = %s LIMIT 1", (doc_id,))
            return cur.fetchone() is not None
        finally:
            conn.close()
    except Exception:
        return False


def _caption_from_parts(parts: list[str]) -> str:
    for part in parts:
        for line in part.splitlines():
            line = line.strip()
            if line and _CAPTION_RE.search(line):
                return line[:300]
    return ""


def parse(path: str | Path) -> ParsedDoc:
    """Extract text from every shape and speaker notes per slide.

    Each slide is prefixed with a '--- Slide N ---' header so the chunker
    produces slide-scoped prose chunks. Embedded tables become Table objects.
    Returns ParsedDoc with text=joined slides and tables=list of Table.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    path = Path(path)
    doc_id = _doc_id_slug(path)
    prs = Presentation(path)

    slides_content: list[dict[str, object]] = []
    tables: list[Table] = []
    images: list[dict] = []
    skip_images = _already_extracted(doc_id)
    image_dir: Path | None = None
    if not skip_images:
        try:
            from config import get_settings

            image_dir = Path(get_settings().image_store_path)
            image_dir.mkdir(parents=True, exist_ok=True)
        except Exception:
            log.warning("pptx_parser: could not create image store; skipping image extraction")
            skip_images = True

    for slide_num, slide in enumerate(prs.slides, 1):
        parts: list[str] = []
        image_idx = 0
        title_text = ""

        try:
            title = getattr(slide.shapes, "title", None)
            title_text = normalize_math_unicode(title.text).strip() if title is not None else ""
            if title_text:
                parts.append(f"Title: {title_text}")
        except Exception:
            pass

        for shape in slide.shapes:
            # Text frames (title, content, text boxes)
            if shape.has_text_frame:
                shape_text = normalize_math_unicode(shape.text_frame.text).strip()
                if title_text and shape_text == title_text:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = normalize_math_unicode(para.text).strip()
                    if line:
                        parts.append(line)

            # Embedded tables
            if shape.has_table:
                rows = [
                    [cell.text.strip() for cell in row.cells]
                    for row in shape.table.rows
                ]
                tables.append(Table(name=f"slide_{slide_num}_table", rows=rows))

            if not skip_images and image_dir is not None and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    if len(blob) < _MIN_IMAGE_BYTES or len(blob) > _MAX_IMAGE_BYTES:
                        continue
                    image_idx += 1
                    ext = str(getattr(shape.image, "ext", None) or "png").lower()
                    filename = f"{_image_filename_slug(path)}_slide{slide_num}_img{image_idx}.{ext}"
                    out_path = image_dir / filename
                    out_path.write_bytes(blob)
                    images.append({
                        "filename": filename,
                        "path": str(out_path),
                        "page": slide_num,
                        "slide": slide_num,
                        "index": image_idx,
                        "caption": _caption_from_parts(parts),
                        "url": f"/images/{filename}",
                    })
                except Exception:
                    log.warning("pptx_parser: image extraction failed for %s slide %s", path.name, slide_num)

        # Speaker notes
        if slide.has_notes_slide:
            notes_text = normalize_math_unicode(slide.notes_slide.notes_text_frame.text).strip()
            if notes_text and len(notes_text) > 20:
                parts.append(f"[Speaker notes]: {notes_text}")

        if parts:
            combined = "\n".join(parts)
            slides_content.append({
                "slide": slide_num,
                "text": combined,
                "char_count": len(combined),
            })

    merged_chunks: list[str] = []
    buffer: list[dict[str, object]] = []
    buffer_chars = 0
    min_slide_chars = 100
    for slide_info in slides_content:
        buffer.append(slide_info)
        buffer_chars += int(slide_info["char_count"])
        if buffer_chars >= min_slide_chars:
            merged_chunks.append("\n\n".join(
                f"--- Slide {s['slide']} ---\n{s['text']}" for s in buffer
            ))
            buffer = []
            buffer_chars = 0
    if buffer:
        merged_chunks.append("\n\n".join(
            f"--- Slide {s['slide']} ---\n{s['text']}" for s in buffer
        ))

    text = normalize_math_unicode("\n\n".join(merged_chunks))

    return ParsedDoc(
        text=text,
        tables=tables,
        metadata={
            "filename": path.name,
            "type": "pptx",
            "slides": len(prs.slides),
        },
        images=images,
        source_ref={"filename": path.name, "slide_count": len(prs.slides)},
    )
